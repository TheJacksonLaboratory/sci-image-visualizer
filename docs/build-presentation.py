#!/usr/bin/env python3
"""Generate the stakeholder-facing jit-ui visualization / jax-image-visualization
presentation as a .pptx file.

Run with:
    python3 libs/jax-image-visualization/docs/build-presentation.py

Produces:
    libs/jax-image-visualization/docs/jax-image-visualization-presentation.pptx

Slide content is the companion deck to JIT_UI_visualization_library_SOW (generated
by gen_jit_ui_visualization_sow.py). Keep the two in sync — when the SOW changes,
update the slide-content here. Mirrors the style of jit-dianne-worker's
docs/build-presentation.py.

Dependencies: python-pptx (`pip install python-pptx`).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
OUTPUT = HERE / "jax-image-visualization-presentation.pptx"
IMG = HERE / "img"
VIS = IMG / "sow-visuals"

# Palette — shared with the DIANNE deck for cross-deck consistency.
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HEADER = RGBColor(0x0B, 0x2B, 0x53)        # deep navy
COLOR_ACCENT = RGBColor(0x1F, 0x7A, 0x8C)        # teal
COLOR_OSD = RGBColor(0x4C, 0x6E, 0xF5)           # blue  — OpenSeadragon
COLOR_PLOTLY = RGBColor(0xE5, 0x73, 0x73)        # warm red — Plotly
COLOR_DELIVERED = RGBColor(0x2E, 0x7D, 0x32)     # green
COLOR_INPROGRESS = RGBColor(0xE6, 0x8A, 0x00)    # amber
COLOR_PLANNED = RGBColor(0x6E, 0x6E, 0x6E)       # grey
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_SECONDARY = RGBColor(0x6E, 0x6E, 0x6E)
COLOR_CHIP_BG = RGBColor(0xEE, 0xF1, 0xF7)

STATUS_COLOR = {
    "Delivered": COLOR_DELIVERED,
    "In progress": COLOR_INPROGRESS,
    "Planned": COLOR_PLANNED,
}

DECK_NAME = "jit-ui Visualization Rework & Library Extraction"


def _enable_wrap(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE


def set_text(frame, text, *, size, bold=False, color=COLOR_BODY, align=None):
    _enable_wrap(frame)
    frame.text = ""
    p = frame.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullets(frame, bullets, *, size=18, color=COLOR_BODY):
    _enable_wrap(frame)
    frame.text = ""
    for i, b in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_title(slide, text, *, subtitle=None, y=0.4):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(y),
                                         Inches(12.5), Inches(0.7))
    set_text(title_box.text_frame, text, size=28, bold=True, color=COLOR_HEADER)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.7),
                                       Inches(12.5), Inches(0.4))
        set_text(sub.text_frame, subtitle, size=16, color=COLOR_SECONDARY)


def add_footer(slide, n, total):
    foot = slide.shapes.add_textbox(Inches(0.6), Inches(7.0),
                                    Inches(12.5), Inches(0.3))
    set_text(foot.text_frame, f"{DECK_NAME}   ·   {n} / {total}",
             size=10, color=COLOR_SECONDARY)


def panel(slide, x, y, w, h, *, color=COLOR_CHIP_BG):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    return box


def add_framed_image(slide, path, *, cx, top, max_w, max_h, caption, cap_y):
    """Add a picture fit within (max_w × max_h), centered on cx, with a caption."""
    pic = slide.shapes.add_picture(str(path), Inches(0), Inches(top),
                                   width=Inches(max_w))
    if pic.height > Inches(max_h):
        scale = Inches(max_h) / pic.height
        pic.width = int(pic.width * scale)
        pic.height = Inches(max_h)
    pic.left = Inches(cx) - pic.width // 2
    pic.top = Inches(top)
    cap = slide.shapes.add_textbox(Inches(cx - max_w / 2), Inches(cap_y),
                                   Inches(max_w), Inches(0.5))
    set_text(cap.text_frame, caption, size=12, color=COLOR_SECONDARY,
             align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- slides ----

def slide_title(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(12.0), Inches(1.6))
    set_text(t.text_frame, "jit-ui Visualization Rework\n& Library Extraction",
             size=46, bold=True, color=COLOR_HEADER)
    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(12.0), Inches(0.8))
    set_text(s.text_frame,
             "A backend-neutral, reusable image-visualization library extracted from jit-ui",
             size=22, color=COLOR_ACCENT)
    m = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(12.0), Inches(0.9))
    set_text(m.text_frame,
             "Companion deck for JIT_UI_visualization_library_SOW   ·   "
             "library: jax-image-visualization\n"
             "Status: substantially delivered and merged to master (PR #79); "
             "library extraction in progress, distribution pending.",
             size=14, color=COLOR_SECONDARY)


def slide_summary(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What this is",
              subtitle="Rework the jit-ui visualization layer, then carve it into a standalone Angular library")
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.8))
    add_bullets(body.text_frame, [
        "One backend-neutral contract (IVisualizer) with two pluggable rendering backends.",
        "OpenSeadragon — natively-tiled, deeply-zoomable Image plot type, off the jit-service tile endpoints.",
        "Plotly — the scientific plot types (heatmap, contour, scatter, surface, scatter-3D, isosurface).",
        "A backend-neutral region model + on-canvas tools that work identically on either backend.",
        "A Channels and Histogram tool: per-channel brightness/contrast, Fiji-style pseudo-colour, true 16-bit.",
        "Extracted into the Nx buildable library jax-image-visualization — the host reaches everything app-specific through DI ports, so the package carries no host coupling.",
    ], size=18)
    add_footer(slide, n, total)


def slide_architecture(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Architecture",
              subtitle="One contract, a router, two backends, one shared state")
    # Left: bullets. Right: the architecture diagram.
    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(6.0), Inches(4.8))
    add_bullets(left.text_frame, [
        "Public API: IVisualizer (IDataRenderer + IRegionStore + IToolController + IDisplayOptions + ViewerCapabilities), IChannelHistogramApi, IRegionEditorApi.",
        "RoutingVisualizerService implements the API tokens and routes Image → OpenSeadragon, everything else → Plotly.",
        "Both backends read one shared VisualizerStore (display / channel / colormap) and one RegionStore (regions, selection, class colours).",
        "Capability gating: consumers call only what the active backend advertises.",
    ], size=15)
    arch = IMG / "jit-ui-visualization-architecture.png"
    if arch.exists():
        add_framed_image(slide, arch, cx=10.0, top=1.6, max_w=5.6, max_h=4.9,
                         caption="jax-image-visualization architecture",
                         cap_y=6.55)
    add_footer(slide, n, total)


def slide_ports(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Host integration — DI ports & adapters",
              subtitle="Everything app-specific is a port; the host provides one adapter each. The library imports nothing from jit-ui.")
    ports = [
        ("VISUALIZER · CHANNEL_HISTOGRAM_API · REGION_EDITOR_API",
         "RoutingVisualizerService (useExisting) — the render/viewport, channel-histogram and region-editor APIs the embeddable components call."),
        ("TILE_ACCESS_PORT",
         "Selected image's base64 FileInfo, bearer auth, ROI-zoom — used to fetch tiles / histograms / exports from jit-service."),
        ("IMAGE_STATE_PORT",
         "Current image info, filename, loading & cache progress, zoom flag, panel width (read + write)."),
        ("REGION_IO_PORT",
         "Selected filename, ROI-file existence check, GeoJSON save — bridges region import/export persistence."),
        ("VIZ_CONFIG",
         "Backend base URL (slideCropServer) and optional UI hints (e.g. regionEditorWidthSelector)."),
    ]
    y = 1.7
    for name, desc in ports:
        panel(slide, 0.6, y, 12.1, 0.86)
        nm = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.08), Inches(11.7), Inches(0.34))
        set_text(nm.text_frame, name, size=14, bold=True, color=COLOR_HEADER)
        ds = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.42), Inches(11.7), Inches(0.40))
        set_text(ds.text_frame, desc, size=12, color=COLOR_BODY)
        y += 1.00
    add_footer(slide, n, total)


def slide_consume(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Consuming the library \u2014 two modes",
              subtitle="Same <visualization> + 4 DI ports; IImageInfo.tiled selects how pixels are sourced")
    col_y, col_h = 1.7, 4.3
    panel(slide, 0.6, col_y, 6.0, col_h)
    ha = slide.shapes.add_textbox(Inches(0.85), Inches(col_y + 0.18), Inches(5.5), Inches(0.4))
    set_text(ha.text_frame, "Mode A \u2014 server-backed", size=18, bold=True, color=COLOR_HEADER)
    ba = slide.shapes.add_textbox(Inches(0.85), Inches(col_y + 0.8), Inches(5.5), Inches(col_h - 1.0))
    add_bullets(ba.text_frame, [
        "Host runs a jit-service-like backend; IMAGE_STATE_PORT emits IImageInfo with tiled: true + urls[] (/preview).",
        "OSD tiles via /tiles/info + /tile off VIZ_CONFIG.slideCropServer; TILE_ACCESS_PORT supplies the info blob, auth, and /zoom/region hi-def re-fetch.",
        "Endpoints: /preview, /tiles/info, /tile (+ /zoom/region, /histogram, /export/tiff for 16-bit).",
        "For whole-slide, deep-zoom, multi-channel 16-bit. (jit-ui uses this.)",
    ], size=13)
    panel(slide, 6.9, col_y, 5.8, col_h)
    hb = slide.shapes.add_textbox(Inches(7.15), Inches(col_y + 0.18), Inches(5.3), Inches(0.4))
    set_text(hb.text_frame, "Mode B \u2014 serverless", size=18, bold=True, color=COLOR_HEADER)
    bb = slide.shapes.add_textbox(Inches(7.15), Inches(col_y + 0.8), Inches(5.3), Inches(col_h - 1.0))
    add_bullets(bb.text_frame, [
        "No backend: host turns in-memory pixels into a blob:/data: URL; IMAGE_STATE_PORT emits tiled: false + urls: [thatUrl].",
        "OSD loads that single image directly \u2014 no /tiles/info, /tile, or slideCropServer. The other ports are no-ops.",
        "Host builds ~no server code \u2014 just the image-state adapter + stubs.",
        "For viewport-sized, decoded images (e.g. the pipeline preview). Trade-off: no deep tiles / hi-def re-fetch / 16-bit.",
    ], size=13)
    note = slide.shapes.add_textbox(Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.5))
    set_text(note.text_frame,
             "Switch: IImageInfo.tiled === false \u2192 direct single-image load (B); otherwise OSD tiles via the server (A). "
             "Same regions / tools / Channels either way.",
             size=12, color=COLOR_SECONDARY)
    add_footer(slide, n, total)


def slide_scope(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Scope")
    col_y, col_h = 1.7, 4.7
    panel(slide, 0.6, col_y, 6.0, col_h)
    h1 = slide.shapes.add_textbox(Inches(0.85), Inches(col_y + 0.18), Inches(5.5), Inches(0.4))
    set_text(h1.text_frame, "In scope", size=20, bold=True, color=COLOR_HEADER)
    b1 = slide.shapes.add_textbox(Inches(0.85), Inches(col_y + 0.8), Inches(5.5), Inches(col_h - 1.0))
    add_bullets(b1.text_frame, [
        "Backend-neutral contract + capability gating.",
        "OpenSeadragon Image backend (tiles, navigator, scale bar, slice cache, recolor).",
        "Plotly scientific plot types + plot-type router.",
        "Region model, tools (wand / vertex-eraser / zoom-to-box / bezier), QuPath GeoJSON I/O.",
        "Channels and Histogram tool incl. true 16-bit + exports.",
        "In-browser SAM segmentation: box / point prompts + automatic Cellpose-SAM (ONNX, pluggable model registry).",
        "Extraction into the Nx buildable library.",
    ], size=13)
    panel(slide, 6.9, col_y, 5.8, col_h)
    h2 = slide.shapes.add_textbox(Inches(7.15), Inches(col_y + 0.18), Inches(5.3), Inches(0.4))
    set_text(h2.text_frame, "Out of scope", size=20, bold=True, color=COLOR_HEADER)
    b2 = slide.shapes.add_textbox(Inches(7.15), Inches(col_y + 0.8), Inches(5.3), Inches(col_h - 1.0))
    add_bullets(b2.text_frame, [
        "The jit-service tile/histogram/export endpoints (their own work, issue #76).",
        "Per-consumer integration of downstream frontends (owned by each consumer).",
        "Server-side rendering or a non-Angular distribution.",
        "Replacing Plotly or OpenSeadragon with a different engine.",
    ], size=13)
    add_footer(slide, n, total)


def _deliverable_col(slide, x, y, w, items):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4.6))
    tf = box.text_frame
    _enable_wrap(tf)
    tf.text = ""
    for i, (code, title, status) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = f"{code} — {title}   "
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.color.rgb = COLOR_HEADER
        s = p.add_run()
        s.text = f"[{status}]"
        s.font.size = Pt(13)
        s.font.bold = True
        s.font.color.rgb = STATUS_COLOR[status]


def slide_deliverables(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Deliverables",
              subtitle="'Delivered' = implemented, verified and merged to master")
    _deliverable_col(slide, 0.6, 1.7, 6.0, [
        ("D1", "Backend-neutral contract & router", "Delivered"),
        ("D2", "OpenSeadragon backend + server tiles", "Delivered"),
        ("D3", "Regions, on-canvas tools, GeoJSON I/O", "Delivered"),
        ("D4", "Plot-type framework (Plotly, profiles)", "Delivered"),
    ])
    _deliverable_col(slide, 6.9, 1.7, 5.8, [
        ("D5", "Channels & Histogram + true 16-bit", "Delivered"),
        ("D6", "Browser SAM segmentation (box/point/Cellpose-SAM)", "Delivered"),
        ("D7", "Extraction into the Nx library", "In progress"),
        ("D8", "Distribution & downstream adoption", "Planned"),
    ])
    add_footer(slide, n, total)


def slide_visuals_a(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Current visuals — imaging & regions")
    add_framed_image(slide, VIS / "fig3-multichannel-fluorescence.png",
                     cx=3.5, top=1.7, max_w=6.0, max_h=4.3,
                     caption="Multi-fluorescence image with 4 channels (Channels & Histogram).",
                     cap_y=6.2)
    add_framed_image(slide, VIS / "fig4-region-bezier.png",
                     cx=9.9, top=1.7, max_w=6.0, max_h=4.3,
                     caption="Region overlay with Bézier curve / polygon shapes.",
                     cap_y=6.2)
    add_footer(slide, n, total)


def slide_visuals_b(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Current visuals — large images, 3D & scientific plots")
    add_framed_image(slide, VIS / "fig5-openseadragon-tiling.png",
                     cx=2.6, top=1.8, max_w=4.1, max_h=3.6,
                     caption="Large whole-slide image via OpenSeadragon & tiling.",
                     cap_y=5.65)
    add_framed_image(slide, VIS / "fig6-ct-isosurface.png",
                     cx=6.85, top=1.8, max_w=4.1, max_h=3.6,
                     caption="CT 3D dataset (nii) — image stack & isosurface.",
                     cap_y=5.65)
    add_framed_image(slide, VIS / "fig7-plotly-contour.png",
                     cx=11.1, top=1.8, max_w=4.1, max_h=3.6,
                     caption="Plotly contour plotting.",
                     cap_y=5.65)
    add_footer(slide, n, total)


def slide_phasing(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phasing & test coverage",
              subtitle="Phases 0–4 delivered · Phase 5 in progress · Phase 6 planned")
    left = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(7.2), Inches(4.6))
    add_bullets(left.text_frame, [
        "Phase 0 — Contract & router  [Delivered]",
        "Phase 1 — OpenSeadragon image backend  [Delivered]",
        "Phase 2 — Regions & tools  [Delivered]",
        "Phase 3 — Plot-type framework  [Delivered]",
        "Phase 4 — Channels & Histogram + 16-bit  [Delivered]",
        "Phase 5 — Library extraction hardening + examples  [In progress]",
        "Phase 6 — Distribution + downstream consumer  [Planned]",
    ], size=16)
    panel(slide, 8.1, 1.7, 4.6, 3.5)
    h = slide.shapes.add_textbox(Inches(8.35), Inches(1.9), Inches(4.1), Inches(0.4))
    set_text(h.text_frame, "Test coverage", size=18, bold=True, color=COLOR_HEADER)
    cov = slide.shapes.add_textbox(Inches(8.35), Inches(2.45), Inches(4.1), Inches(2.6))
    add_bullets(cov.text_frame, [
        "Jest suite — 682 tests / 43 suites, all passing — in CI with enforced minimum thresholds.",
        "Lines 73% · Statements 71% · Functions 63% · Branches 56%.",
        "Build fails if coverage regresses below the gate.",
    ], size=13)
    add_footer(slide, n, total)


def slide_risks(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Risks & mitigations")
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.8))
    add_bullets(body.text_frame, [
        "Residual host coupling blocks a clean extraction → dependency-boundary lint rule; route every app-specific behaviour through a DI port.",
        "ng-packagr / AOT constraints (NG3001, partial compilation) → keep the public API in the barrel; build the library in CI on every change.",
        "OpenSeadragon tiling/recolor perf on large grayscale stacks → precomputed recolor LUT (done); bound per-channel tile counts.",
        "16-bit display fidelity limited (viewer tiles are 8-bit) → precision preserved in the native histogram and the 16-bit TIFF export.",
        "Versioning / breaking changes destabilize consumers → semantic versioning + documented breaking-change policy (D8); small contract-first surface.",
    ], size=16)
    add_footer(slide, n, total)


def _stat_panel(slide, x, y, w, h, *, label, big, sub, big_color):
    """A rounded callout: small label, big number, small sub-caption."""
    panel(slide, x, y, w, h)
    tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12),
                                  Inches(w - 0.3), Inches(h - 0.2))
    tf = tb.text_frame
    _enable_wrap(tf)
    tf.text = ""
    rows = [(label, 12, False, COLOR_SECONDARY),
            (big, 32, True, big_color),
            (sub, 11, False, COLOR_SECONDARY)]
    for i, (text, size, bold, color) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def slide_dedicated_repo(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Recommendation: a dedicated repo, not the jds-ui-toolkit monorepo",
              subtitle="jax-image-visualization differs from the toolkit's libs in scale, dependencies, and audience")

    # The size mismatch, up top — three stat callouts.
    w = 3.95
    gap = 0.35
    _stat_panel(slide, 0.6, 1.55, w, 1.4,
                label="jax-image-visualization",
                big="≈20,600", sub="lines (src, excl. tests)", big_color=COLOR_ACCENT)
    _stat_panel(slide, 0.6 + (w + gap), 1.55, w, 1.4,
                label="jds-ui-toolkit — all libs",
                big="≈5,500", sub="components 2,963 · api-clients 2,146 · themes 380", big_color=COLOR_SECONDARY)
    _stat_panel(slide, 0.6 + 2 * (w + gap), 1.55, w, 1.4,
                label="relative size",
                big="≈3.8×", sub="the whole toolkit  ·  ~7× components (its largest)", big_color=COLOR_OSD)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(12.1), Inches(3.6))
    add_bullets(body.text_frame, [
        "Scale & identity — at ~20.6k LOC the viz lib is ~3.8× all three jds-ui-toolkit libs combined (components 2,963 + api-clients 2,146 + themes 380 ≈ 5,500) and ~7× the largest, components. It would dominate and reshape a toolkit built for small, shared UI pieces.",
        "Heavy, specialized dependencies — OpenSeadragon, Plotly, image-js, onnxruntime-web (in-browser SAM/Cellpose segmentation), and optional cellpose-js. These large WebGL/WASM deps inflate install size, CI time, and the dependency surface for consumers that only want a button or a theme.",
        "Independent release cadence — fast-moving features (OSD, regions, segmentation) need their own semver + breaking-change policy; a shared repo forces lock-step releases or per-package publishing that churns consumers of tiny libs.",
        "External contributors / OSS readiness — a dedicated repo exposes only the publishable library (its own LICENSE, README, CHANGELOG, scoped issues/PRs, narrow access). A shared monorepo would expose unrelated internal libs (api-clients, internal components) — wider IP/security surface and CODEOWNERS complexity.",
        "Focused CI & lean clones — build/test scoped to one package, no cross-package graph to configure; large doc/image assets (this deck alone is ~15 MB) stay out of every toolkit clone.",
        "Clear ownership & decoupled consumption — consumers depend on a pinned, published npm package, not a workspace path. Monorepos pay off for tightly-coupled, co-released packages; this is a standalone, contract-first, reusable library with a different audience.",
    ], size=13)
    add_footer(slide, n, total)


def slide_distribution(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Distribution & open questions",
              subtitle="npm under the @jax-data-science org scope (Nx publishable library)")
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.8))
    add_bullets(body.text_frame, [
        "Package name TBD: @jax-data-science/{image-visualization | jax-image-visualization | jit-image-visualization}.",
        "Public npm vs a private/internal registry — confirm the Nx publishable-library release setup.",
        "Theming / customization API — how do consumers override colours, tools, and the tile-source backend?",
        "Which residual components (if any) are too jit-ui-coupled to move yet — and the deferral plan.",
        "First downstream consumer — which JAX Data Science frontend proves reuse, and on what timeline?",
    ], size=16)
    add_footer(slide, n, total)


def slide_future_work(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Future work & investigations",
              subtitle="Exploratory — not part of the committed deliverables (§3)")
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.8))
    add_bullets(body.text_frame, [
        "More segmentation models (extends D6) — additional fine-tuned Cellpose-SAM variants for other tissue/cell domains, and further promptable models (e.g. SAM3). They drop into the existing pluggable registry, so adding one is largely export-and-host.",
        "3D / z-stack segmentation — propagate SAM masks across slices with cross-slice linking, so a 2D prompt extends through a volume.",
        "Image-visualization MCP server — expose the viewer as Model Context Protocol tools/resources so an LLM client can drive it programmatically (open an image, switch plot type, pan/zoom to a region, toggle channels/colormap, create/select/edit regions, read back the current view or region geometry).",
        "Goal of the MCP server: agentic / natural-language control (e.g. “zoom to the largest region and export it as GeoJSON”) and headless automation. Builds on the stable library contract (D7); transport, scope, and security/auth boundaries TBD.",
    ], size=15)
    add_footer(slide, n, total)


def slide_mcp_server(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Future work \u2014 Image-visualization MCP server",
              subtitle="Drive the live viewer from a Claude / MCP client \u2014 a control bridge, not a renderer (jit-ui#97)")
    # Control-flow strip.
    panel(slide, 0.6, 1.55, 12.1, 0.8)
    flow = slide.shapes.add_textbox(Inches(0.8), Inches(1.72), Inches(11.7), Inches(0.5))
    set_text(flow.text_frame,
             "Claude / MCP client  \u2192  jax-image-viz-mcp (Node, MCP stdio/SSE)  \u2192  \u2026  \u2192  "
             "RoutingVisualizerService  (live browser tab)",
             size=13, bold=True, color=COLOR_HEADER, align=PP_ALIGN.CENTER)
    body = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.1), Inches(4.0))
    add_bullets(body.text_frame, [
        "Not a stateless plot generator \u2014 this is a stateful, browser-resident viewer (OpenSeadragon + Plotly). \u201cControl the plotting\u201d means driving the live UI, not rendering a file.",
        "The control surface already exists \u2014 the whole UI runs through one backend-neutral IVisualizer contract (~80 methods); the server exposes that, it doesn't invent a new API.",
        "Tools generated from one command catalogue \u2014 viz.loadImage, viz.setPlotType, viz.setColormap, viz.setRegionsGeoJson, viz.segmentRectangles[Cellpose], viz.getScreenshotPng\u2026 \u2014 capability-gated so Claude is never offered a no-op tool.",
        "MVP: getState + 5 commands (load \u2192 setPlotType \u2192 setColormap \u2192 setRegionsGeoJson \u2192 getScreenshot) proves the transport / handshake / serializable boundary before expanding the catalogue.",
        "Security: the channel sits behind oauth2-proxy + Auth0 with a viewerId handshake; ships dev/flagged first, hardened before any controlled deployment.",
    ], size=13)
    add_footer(slide, n, total)


def slide_close(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(12.0), Inches(1.6))
    set_text(t.text_frame, "One contract. Two backends. A reusable library.",
             size=40, bold=True, color=COLOR_HEADER)
    s = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(12.0), Inches(1.4))
    add_bullets(s.text_frame, [
        "Source: jit-ui master (merged via PR #79) — libs/jax-image-visualization/.",
        "Detailed SOW: JIT_UI_visualization_library_SOW (gen_jit_ui_visualization_sow.py).",
        "Next: finish the extraction (Phase 5), then publish and prove reuse in a second frontend (Phase 6).",
    ], size=16, color=COLOR_BODY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_makers = [
        slide_title,
        slide_summary,
        slide_architecture,
        slide_ports,
        slide_consume,
        slide_scope,
        slide_deliverables,
        slide_visuals_a,
        slide_visuals_b,
        slide_phasing,
        slide_risks,
        slide_dedicated_repo,
        slide_distribution,
        slide_future_work,
        slide_mcp_server,
        slide_close,
    ]
    total = len(slide_makers)
    for i, make in enumerate(slide_makers, start=1):
        if make is slide_title:
            make(prs, total)
        else:
            make(prs, i, total)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build()
